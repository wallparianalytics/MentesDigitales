#!/usr/bin/env python3
"""
Script para corregir los formatos de Sesión 6 después de crearla
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def corregir_formatos():
    """Corrige los formatos de la presentación"""
    prs = Presentation("Sesión 6.pptx")

    # Colores del tema
    COLOR_AZUL_HEADER = RGBColor(99, 102, 241)
    COLOR_GRIS_OSCURO = RGBColor(30, 41, 59)
    COLOR_BLANCO = RGBColor(255, 255, 255)

    # Mapa de tamaños de fuente por tipo de contenido
    formato_por_contenido = {
        'header': {'fuente': 'Arial', 'tamaño': 27, 'color': COLOR_BLANCO, 'negrita': True},
        'titulo_emoji': {'fuente': 'Arial', 'tamaño': 18, 'color': COLOR_AZUL_HEADER, 'negrita': True},
        'titulo_normal': {'fuente': 'Arial', 'tamaño': 15, 'color': COLOR_AZUL_HEADER, 'negrita': True},
        'texto_normal': {'fuente': 'Arial', 'tamaño': 13.5, 'color': COLOR_GRIS_OSCURO, 'negrita': False},
        'texto_grande': {'fuente': 'Arial', 'tamaño': 18, 'color': COLOR_BLANCO, 'negrita': False},
        'emoji_grande': {'fuente': 'Arial', 'tamaño': 36, 'color': COLOR_GRIS_OSCURO, 'negrita': False},
    }

    # Procesar todas las diapositivas excepto la primera (portada)
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue  # Saltar portada que ya está bien

        print(f"Procesando diapositiva {idx + 1}...")

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            texto = shape.text.strip()
            if not texto:
                continue

            # Determinar el formato basado en el contenido
            formato = None
            if any(emoji in texto for emoji in ['🎯', '📋', '💭', '⏱️', '📝', '💬', '📌', '🚀', '⚠️']):
                if len(texto) < 20:
                    formato = 'titulo_emoji'
                else:
                    formato = 'texto_normal'
            elif len(texto) == 1 or len(texto) == 2:  # Emojis solos
                formato = 'emoji_grande'
            elif shape.top < 1000000:  # Header (top < 1 inch aproximadamente)
                formato = 'header'
            elif len(texto) > 50:
                # Texto largo, probablemente descripción
                # Si está en posición alta (probablemente sobre imagen de fondo), usar texto grande blanco
                if shape.top < 3500000:  # Top < 3.5 inches aproximadamente
                    formato = 'texto_grande'
                else:
                    formato = 'texto_normal'
            else:
                formato = 'texto_normal'

            if formato and formato in formato_por_contenido:
                config = formato_por_contenido[formato]

                # Aplicar formato a todos los runs
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = config['fuente']
                        run.font.size = Pt(config['tamaño'])
                        run.font.color.rgb = config['color']
                        if config['negrita']:
                            run.font.bold = True

    # Guardar la presentación corregida
    prs.save("Sesión 6.pptx")
    print("\n✅ Formatos corregidos exitosamente!")

if __name__ == "__main__":
    corregir_formatos()
