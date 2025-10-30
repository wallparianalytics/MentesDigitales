#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import copy

def crear_sesion_6():
    # Cargar la presentación de la Sesión 5 como plantilla
    prs = Presentation('/home/user/MentesDigitales/Sesión 5.pptx')

    # Colores del tema (basados en Sesión 5)
    COLOR_AZUL_HEADER = RGBColor(99, 102, 241)  # Azul índigo para headers
    COLOR_NARANJA = RGBColor(245, 158, 11)  # Naranja/amarillo para decoración
    COLOR_GRIS_CLARO = RGBColor(241, 245, 249)  # Gris claro para cajas
    COLOR_BLANCO = RGBColor(255, 255, 255)  # Blanco

    # Mantener solo la primera diapositiva (portada) y eliminar el resto
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]

    # Modificar la portada (Diapositiva 1)
    slide1 = prs.slides[0]

    # Actualizar textos de la portada
    for shape in slide1.shapes:
        if hasattr(shape, "text"):
            if "Sesión 5" in shape.text:
                shape.text = "Sesión 6"
            elif "Prototipar tu Proyecto de IA" in shape.text:
                shape.text = "Recopilación y Preparación de Datos"
            elif "Octubre 2024" in shape.text:
                shape.text = "Noviembre 2024"

    # Función auxiliar para crear una diapositiva con header azul
    def crear_slide_con_header(titulo):
        # Usar el primer layout disponible
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Eliminar todas las formas del layout
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)

        # Rectángulo azul superior
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.96)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_AZUL_HEADER
        header.line.fill.background()

        # Línea delgada debajo del header
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.94),
            Inches(10), Inches(0.02)
        )
        line.line.fill.background()

        # Título en el header
        title_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.25),
            Inches(7), Inches(0.42)
        )
        title_frame = title_box.text_frame
        title_frame.text = titulo
        p = title_frame.paragraphs[0]
        p.font.size = Pt(27)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.LEFT

        return slide

    # Función para agregar logo UPC (copiar de la slide 1)
    def agregar_logo(slide):
        # Buscar el logo en la portada
        for shape in slide1.shapes:
            if shape.shape_type == 13:  # PICTURE
                if shape.left.inches > 8:  # Logo en esquina superior derecha
                    # Copiar el logo
                    el = shape.element
                    newel = copy.deepcopy(el)
                    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                    break

    # DIAPOSITIVA 2: Objetivo y Agenda
    slide2 = crear_slide_con_header("Objetivo y Agenda")
    agregar_logo(slide2)

    # Caja de objetivo
    obj_box = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.50), Inches(1.65),
        Inches(9.00), Inches(1.23)
    )
    obj_box.fill.solid()
    obj_box.fill.fore_color.rgb = COLOR_GRIS_CLARO
    obj_box.line.fill.background()

    # Línea vertical en la caja
    line_vert = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.53), Inches(1.65),
        Inches(0.03), Inches(1.23)
    )
    line_vert.line.fill.background()

    # Título del objetivo
    obj_title = slide2.shapes.add_textbox(
        Inches(0.81), Inches(1.90),
        Inches(4.31), Inches(0.33)
    )
    tf = obj_title.text_frame
    tf.text = "🎯 Objetivo"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Texto del objetivo
    obj_text = slide2.shapes.add_textbox(
        Inches(0.81), Inches(2.30),
        Inches(8.61), Inches(0.50)
    )
    tf = obj_text.text_frame
    tf.word_wrap = True
    tf.text = "Enseñar la importancia de los datos para entrenar modelos de Inteligencia Artificial y cómo recopilar datos de calidad."
    p = tf.paragraphs[0]
    p.font.size = Pt(13.5)
    p.alignment = PP_ALIGN.LEFT

    # Título de agenda
    agenda_title = slide2.shapes.add_textbox(
        Inches(0.50), Inches(3.13),
        Inches(9.18), Inches(0.29)
    )
    tf = agenda_title.text_frame
    tf.text = "📋 Agenda"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Contenido de agenda
    agenda_text = slide2.shapes.add_textbox(
        Inches(0.50), Inches(3.56),
        Inches(9.00), Inches(1.29)
    )
    tf = agenda_text.text_frame
    tf.word_wrap = True

    items = [
        "Importancia de los datos en IA",
        "Tipos de datos para entrenar modelos",
        "Cómo recopilar datos de calidad",
        "Actividad práctica: recopilación de datos"
    ]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13.5)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 3: ¿Por qué son importantes los datos?
    slide3 = crear_slide_con_header("¿Por qué son importantes los datos?")
    agregar_logo(slide3)

    # Texto introductorio
    intro_box = slide3.shapes.add_textbox(
        Inches(0.58), Inches(1.40),
        Inches(8.84), Inches(0.80)
    )
    tf = intro_box.text_frame
    tf.word_wrap = True
    tf.text = "Los datos son el combustible de la Inteligencia Artificial. Sin datos de calidad, los modelos de IA no pueden aprender correctamente."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Tres cajas con conceptos clave
    conceptos = [
        ("💡", "Datos = Conocimiento", 0.33),
        ("🎯", "Calidad > Cantidad", 3.61),
        ("📊", "Variedad es clave", 6.89)
    ]

    for emoji, texto, left_pos in conceptos:
        # Caja de fondo
        caja = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(2.60),
            Inches(2.78), Inches(1.50)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_GRIS_CLARO
        caja.line.fill.background()

        # Emoji
        emoji_box = slide3.shapes.add_textbox(
            Inches(left_pos + 1.07), Inches(2.77),
            Inches(0.64), Inches(0.50)
        )
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

        # Texto
        text_box = slide3.shapes.add_textbox(
            Inches(left_pos + 0.50), Inches(3.60),
            Inches(1.78), Inches(0.29)
        )
        tf = text_box.text_frame
        tf.text = texto
        p = tf.paragraphs[0]
        p.font.size = Pt(13.5)
        p.alignment = PP_ALIGN.CENTER

    # DIAPOSITIVA 4: Tipos de Datos para Entrenar IA
    slide4 = crear_slide_con_header("Tipos de Datos para Entrenar IA")
    agregar_logo(slide4)

    # Cuatro cajas con tipos de datos
    tipos = [
        ("📸", "Imágenes", "Fotos de objetos, rostros,\nanimales, plantas", 0.60, 1.50),
        ("🔊", "Sonidos", "Voces, música,\nsonidos ambientales", 5.15, 1.50),
        ("🤸", "Poses", "Movimientos corporales,\nposturas, gestos", 0.60, 3.50),
        ("📝", "Texto", "Palabras, frases,\netiquetas, descripciones", 5.15, 3.50)
    ]

    for emoji, titulo, descripcion, left_pos, top_pos in tipos:
        # Caja de fondo
        caja = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.10), Inches(1.30)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_GRIS_CLARO
        caja.line.fill.background()

        # Emoji y título en la misma línea
        header_box = slide4.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 0.15),
            Inches(3.70), Inches(0.40)
        )
        tf = header_box.text_frame
        tf.text = f"{emoji}  {titulo}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide4.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 0.60),
            Inches(3.70), Inches(0.60)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 5: ¿Cómo Recopilar Datos de Calidad?
    slide5 = crear_slide_con_header("¿Cómo Recopilar Datos de Calidad?")
    agregar_logo(slide5)

    # Contenido con reglas
    content_box = slide5.shapes.add_textbox(
        Inches(0.80), Inches(1.50),
        Inches(8.40), Inches(3.80)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    reglas = [
        ("1️⃣ Cantidad suficiente", "Necesitas muchos ejemplos de cada categoría (mínimo 20-30)"),
        ("2️⃣ Variedad", "Incluye diferentes ángulos, iluminación, fondos, contextos"),
        ("3️⃣ Claridad", "Asegúrate de que los datos sean claros y representativos"),
        ("4️⃣ Balance", "Ten una cantidad similar de datos para cada categoría"),
        ("5️⃣ Relevancia", "Los datos deben relacionarse directamente con tu proyecto")
    ]

    for i, (numero, descripcion) in enumerate(reglas):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"{numero} {descripcion}"
        p.font.size = Pt(15)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(15)
        # Hacer el número en negrita
        if len(p.runs) > 0:
            p.runs[0].font.bold = True

    # DIAPOSITIVA 6: Ejemplo Práctico
    slide6 = crear_slide_con_header("Ejemplo Práctico: Clasificador de Animales")
    agregar_logo(slide6)

    # Caja de título del proyecto
    proyecto_box = slide6.shapes.add_textbox(
        Inches(0.80), Inches(1.40),
        Inches(8.40), Inches(0.50)
    )
    tf = proyecto_box.text_frame
    tf.text = "🎯 Proyecto: Una IA que identifique perros, gatos y pájaros"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Contenido del ejemplo
    ejemplo_box = slide6.shapes.add_textbox(
        Inches(1.20), Inches(2.20),
        Inches(7.60), Inches(3.00)
    )
    tf = ejemplo_box.text_frame
    tf.word_wrap = True

    # Título "Datos necesarios:"
    p = tf.paragraphs[0]
    p.text = "Datos necesarios:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(12)

    ejemplos = [
        "✓ 30+ fotos de perros (diferentes razas, tamaños, colores)",
        "✓ 30+ fotos de gatos (diferentes poses, lugares)",
        "✓ 30+ fotos de pájaros (en vuelo, posados, diferentes especies)",
        "✓ Buena iluminación en todas las fotos",
        "✓ Fondos variados para cada categoría",
        "✓ Diferentes distancias: cerca y lejos"
    ]

    for ejemplo in ejemplos:
        p = tf.add_paragraph()
        p.text = ejemplo
        p.font.size = Pt(14)
        p.space_after = Pt(10)
        p.level = 1

    # DIAPOSITIVA 7: Herramientas para Recopilar Datos
    slide7 = crear_slide_con_header("Herramientas para Recopilar Datos")
    agregar_logo(slide7)

    # Herramientas
    herramientas = [
        ("📱", "Cámara del celular", "Captura imágenes y videos fácilmente", 0.60, 1.50),
        ("🎤", "Micrófono", "Graba sonidos y audio", 5.15, 1.50),
        ("💻", "Teachable Machine", "Captura datos directamente desde el navegador", 0.60, 3.20),
        ("👥", "Colaboración", "Trabaja con compañeros para más variedad", 5.15, 3.20)
    ]

    for emoji, titulo, descripcion, left_pos, top_pos in herramientas:
        # Caja de fondo
        caja = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.10), Inches(1.20)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_GRIS_CLARO
        caja.line.fill.background()

        # Emoji y título
        header_box = slide7.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 0.15),
            Inches(3.70), Inches(0.35)
        )
        tf = header_box.text_frame
        tf.text = f"{emoji}  {titulo}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide7.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 0.55),
            Inches(3.70), Inches(0.55)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 8: Preparando tus Datos
    slide8 = crear_slide_con_header("Preparando tus Datos")
    agregar_logo(slide8)

    # Título de pasos
    pasos_title = slide8.shapes.add_textbox(
        Inches(0.80), Inches(1.40),
        Inches(8.40), Inches(0.40)
    )
    tf = pasos_title.text_frame
    tf.text = "Pasos para preparar tus datos:"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True

    # Lista de pasos
    pasos_box = slide8.shapes.add_textbox(
        Inches(1.20), Inches(2.00),
        Inches(7.60), Inches(3.20)
    )
    tf = pasos_box.text_frame
    tf.word_wrap = True

    pasos = [
        "1. Organiza tus datos por categorías (crea carpetas o clases)",
        "2. Elimina datos borrosos o de mala calidad",
        "3. Revisa que todas las categorías tengan cantidad similar",
        "4. Etiqueta correctamente cada grupo de datos",
        "5. Verifica que los datos representen bien lo que quieres enseñar"
    ]

    for i, paso in enumerate(pasos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = paso
        p.font.size = Pt(16)
        p.space_after = Pt(18)

    # DIAPOSITIVA 9: Actividad Práctica
    slide9 = crear_slide_con_header("🎯 Actividad Práctica")
    agregar_logo(slide9)

    # Título de actividad
    act_title = slide9.shapes.add_textbox(
        Inches(0.80), Inches(1.30),
        Inches(8.40), Inches(0.50)
    )
    tf = act_title.text_frame
    tf.text = "¡Es hora de recopilar datos para tu proyecto!"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Instrucciones
    instr_box = slide9.shapes.add_textbox(
        Inches(1.00), Inches(2.00),
        Inches(8.00), Inches(3.30)
    )
    tf = instr_box.text_frame
    tf.word_wrap = True

    # Título instrucciones
    p = tf.paragraphs[0]
    p.text = "Instrucciones:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(12)

    instrucciones = [
        "1. Revisa el proyecto que planeaste en la Sesión 5",
        "2. Define qué tipo de datos necesitas (imágenes, sonidos, poses)",
        "3. Determina cuántas categorías o clases tendrás",
        "4. Recopila al menos 25-30 ejemplos de cada categoría",
        "5. Organiza tus datos en carpetas o grupos",
        "6. Revisa la calidad y variedad de tus datos"
    ]

    for instruccion in instrucciones:
        p = tf.add_paragraph()
        p.text = instruccion
        p.font.size = Pt(14)
        p.space_after = Pt(10)

    # DIAPOSITIVA 10: Consejos Finales
    slide10 = crear_slide_con_header("💡 Consejos Finales")
    agregar_logo(slide10)

    # Consejos
    consejos_box = slide10.shapes.add_textbox(
        Inches(1.00), Inches(1.50),
        Inches(8.00), Inches(3.80)
    )
    tf = consejos_box.text_frame
    tf.word_wrap = True

    consejos = [
        "✨ Sé creativo: Usa diferentes ángulos, distancias y perspectivas",
        "🎨 Varía el contexto: Diferentes fondos y entornos",
        "☀️ Iluminación: Captura datos con buena luz",
        "🔄 Diversidad: Incluye diferentes variaciones de cada categoría",
        "🤝 Colabora: Trabaja con compañeros para obtener más datos",
        "⏰ Ten paciencia: La recopilación lleva tiempo, ¡pero vale la pena!",
        "🎯 Calidad > Cantidad: Mejor pocos datos buenos que muchos malos"
    ]

    for i, consejo in enumerate(consejos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = consejo
        p.font.size = Pt(15)
        p.space_after = Pt(12)

    # DIAPOSITIVA 11: Próxima Sesión
    slide11 = crear_slide_con_header("Próxima Sesión")
    agregar_logo(slide11)

    # Título de próxima sesión
    prox_title = slide11.shapes.add_textbox(
        Inches(0.50), Inches(2.00),
        Inches(9.00), Inches(0.60)
    )
    tf = prox_title.text_frame
    tf.text = "Sesión 7: Entrenando el Modelo de IA"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Descripción
    prox_desc = slide11.shapes.add_textbox(
        Inches(1.00), Inches(3.00),
        Inches(8.00), Inches(1.00)
    )
    tf = prox_desc.text_frame
    tf.word_wrap = True
    tf.text = "¡Usaremos los datos que recopilaste hoy para entrenar tu modelo de Inteligencia Artificial!"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Guardar presentación
    prs.save('/home/user/MentesDigitales/Sesión 6.pptx')
    print("✅ Presentación 'Sesión 6.pptx' creada exitosamente con el formato de Sesión 5!")

if __name__ == "__main__":
    crear_sesion_6()
