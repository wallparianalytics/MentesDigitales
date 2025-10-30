#!/usr/bin/env python3
"""
Script para comparar los formatos entre Sesión 5 y Sesión 6
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def analizar_presentacion(pptx_path):
    """Analiza una presentación y retorna información de formato"""
    prs = Presentation(pptx_path)

    info = {
        'num_slides': len(prs.slides),
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height,
        'slides': []
    }

    for idx, slide in enumerate(prs.slides):
        slide_info = {
            'numero': idx + 1,
            'layout': slide.slide_layout.name,
            'formas': []
        }

        for shape in slide.shapes:
            shape_info = {
                'tipo': shape.shape_type,
                'nombre': shape.name if hasattr(shape, 'name') else 'Sin nombre',
                'posicion': f"({shape.left}, {shape.top})",
                'tamaño': f"({shape.width}, {shape.height})"
            }

            # Si tiene texto
            if hasattr(shape, "text_frame"):
                shape_info['tiene_texto'] = True
                shape_info['texto'] = shape.text[:50] if shape.text else ""

                # Analizar formato de texto
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        shape_info['fuente'] = {
                            'nombre': run.font.name,
                            'tamaño': run.font.size.pt if run.font.size else "No definido",
                            'negrita': run.font.bold,
                            'cursiva': run.font.italic,
                        }

                        # Color de fuente
                        if run.font.color.type == 1:  # RGB
                            shape_info['fuente']['color'] = f"RGB({run.font.color.rgb[0]}, {run.font.color.rgb[1]}, {run.font.color.rgb[2]})"

                    # Alineación
                    shape_info['alineacion'] = str(para.alignment)

            # Color de fondo
            if hasattr(shape, 'fill'):
                if shape.fill.type == 1:  # SOLID
                    try:
                        rgb = shape.fill.fore_color.rgb
                        shape_info['color_fondo'] = f"RGB({rgb[0]}, {rgb[1]}, {rgb[2]})"
                    except:
                        pass

            slide_info['formas'].append(shape_info)

        info['slides'].append(slide_info)

    return info

def imprimir_comparacion(info5, info6):
    """Imprime una comparación detallada"""
    print("=" * 80)
    print("COMPARACIÓN SESIÓN 5 vs SESIÓN 6")
    print("=" * 80)

    print(f"\nDimensiones de diapositivas:")
    print(f"  Sesión 5: {info5['slide_width']} x {info5['slide_height']}")
    print(f"  Sesión 6: {info6['slide_width']} x {info6['slide_height']}")

    print(f"\nNúmero de diapositivas:")
    print(f"  Sesión 5: {info5['num_slides']}")
    print(f"  Sesión 6: {info6['num_slides']}")

    print("\n" + "=" * 80)
    print("ANÁLISIS DETALLADO POR DIAPOSITIVA")
    print("=" * 80)

    for idx in range(min(3, len(info5['slides']))):  # Primeras 3 diapositivas
        slide5 = info5['slides'][idx]
        print(f"\n--- DIAPOSITIVA {idx + 1} (SESIÓN 5) ---")
        print(f"Layout: {slide5['layout']}")

        for forma in slide5['formas']:
            if forma.get('tiene_texto'):
                print(f"\nForma: {forma['nombre']}")
                print(f"  Texto: {forma['texto']}")
                print(f"  Posición: {forma['posicion']}")
                print(f"  Tamaño: {forma['tamaño']}")
                if 'fuente' in forma:
                    print(f"  Fuente: {forma['fuente']}")
                if 'alineacion' in forma:
                    print(f"  Alineación: {forma['alineacion']}")
                if 'color_fondo' in forma:
                    print(f"  Color fondo: {forma['color_fondo']}")

    print("\n" + "=" * 80)

    for idx in range(min(3, len(info6['slides']))):  # Primeras 3 diapositivas
        slide6 = info6['slides'][idx]
        print(f"\n--- DIAPOSITIVA {idx + 1} (SESIÓN 6) ---")
        print(f"Layout: {slide6['layout']}")

        for forma in slide6['formas']:
            if forma.get('tiene_texto'):
                print(f"\nForma: {forma['nombre']}")
                print(f"  Texto: {forma['texto']}")
                print(f"  Posición: {forma['posicion']}")
                print(f"  Tamaño: {forma['tamaño']}")
                if 'fuente' in forma:
                    print(f"  Fuente: {forma['fuente']}")
                if 'alineacion' in forma:
                    print(f"  Alineación: {forma['alineacion']}")
                if 'color_fondo' in forma:
                    print(f"  Color fondo: {forma['color_fondo']}")

if __name__ == "__main__":
    print("Analizando Sesión 5...")
    info5 = analizar_presentacion("Sesión 5.pptx")

    print("Analizando Sesión 6...")
    info6 = analizar_presentacion("Sesión 6.pptx")

    imprimir_comparacion(info5, info6)
