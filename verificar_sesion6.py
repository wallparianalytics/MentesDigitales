#!/usr/bin/env python3
"""
Script para verificar los formatos de Sesión 6
"""

from pptx import Presentation

def verificar_sesion6():
    """Verifica los formatos de Sesión 6"""
    prs = Presentation("Sesión 6.pptx")

    print("=" * 80)
    print("VERIFICACIÓN DE FORMATOS - SESIÓN 6")
    print("=" * 80)

    # Verificar diapositiva 1 (portada)
    slide1 = prs.slides[0]
    print("\n--- DIAPOSITIVA 1 (PORTADA) ---")
    for shape in slide1.shapes:
        if hasattr(shape, "text_frame") and shape.text:
            print(f"\nTexto: {shape.text}")
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    print(f"  Fuente: {run.font.name}")
                    print(f"  Tamaño: {run.font.size.pt if run.font.size else 'No definido'}pt")
                    print(f"  Negrita: {run.font.bold}")
                    if run.font.color.type == 1:  # RGB
                        rgb = run.font.color.rgb
                        print(f"  Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")

    # Verificar diapositiva 2
    if len(prs.slides) > 1:
        slide2 = prs.slides[1]
        print("\n--- DIAPOSITIVA 2 (OBJETIVO Y AGENDA) ---")
        for shape in slide2.shapes:
            if hasattr(shape, "text_frame") and shape.text and len(shape.text.strip()) > 0:
                texto_preview = shape.text[:50]
                print(f"\nTexto: {texto_preview}")
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        print(f"  Fuente: {run.font.name}")
                        print(f"  Tamaño: {run.font.size.pt if run.font.size else 'No definido'}pt")
                        print(f"  Negrita: {run.font.bold}")
                        if run.font.color.type == 1:  # RGB
                            rgb = run.font.color.rgb
                            print(f"  Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")

    # Verificar diapositiva 3
    if len(prs.slides) > 2:
        slide3 = prs.slides[2]
        print("\n--- DIAPOSITIVA 3 (¿POR QUÉ SON IMPORTANTES LOS DATOS?) ---")
        for shape in slide3.shapes:
            if hasattr(shape, "text_frame") and shape.text and len(shape.text.strip()) > 0:
                texto_preview = shape.text[:50]
                print(f"\nTexto: {texto_preview}")
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        print(f"  Fuente: {run.font.name}")
                        print(f"  Tamaño: {run.font.size.pt if run.font.size else 'No definido'}pt")
                        print(f"  Negrita: {run.font.bold}")
                        if run.font.color.type == 1:  # RGB
                            rgb = run.font.color.rgb
                            print(f"  Color: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")

    print("\n" + "=" * 80)
    print("VERIFICACIÓN COMPLETA")
    print("=" * 80)

if __name__ == "__main__":
    verificar_sesion6()
