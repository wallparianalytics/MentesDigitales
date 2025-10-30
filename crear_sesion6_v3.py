#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import copy

def crear_sesion_6():
    # Cargar la presentación de la Sesión 5 como plantilla
    prs = Presentation('/home/user/MentesDigitales/Sesión 5.pptx')

    # Colores del tema (basados en Sesión 5)
    COLOR_AZUL_HEADER = RGBColor(99, 102, 241)  # Azul índigo para headers
    COLOR_NARANJA = RGBColor(245, 158, 11)  # Naranja/amarillo para decoración
    COLOR_GRIS_CLARO = RGBColor(241, 245, 249)  # Gris claro para cajas
    COLOR_BLANCO = RGBColor(255, 255, 255)  # Blanco
    COLOR_NEGRO = RGBColor(0, 0, 0)  # Negro para texto

    # Mantener solo la primera diapositiva (portada) y eliminar el resto
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]

    # Modificar la portada (Diapositiva 1)
    slide1 = prs.slides[0]

    # Actualizar textos de la portada
    for shape in slide1.shapes:
        if hasattr(shape, "text"):
            if "Sesión 5" in shape.text:
                shape.text = "Sesión 6"
            elif "Prototipar tu Proyecto de IA" in shape.text:
                shape.text = "Recopilación y Preparación de Datos"
            elif "Octubre 2024" in shape.text:
                shape.text = "Noviembre 2024"

    # Función auxiliar para crear una diapositiva con header azul
    def crear_slide_con_header(titulo):
        # Usar el primer layout disponible
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Eliminar todas las formas del layout
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)

        # Rectángulo azul superior
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.96)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_AZUL_HEADER
        header.line.fill.background()

        # Línea delgada debajo del header (casi invisible)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.94),
            Inches(10), Inches(0.00)
        )
        line.line.fill.background()

        # Título en el header
        title_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.25),
            Inches(7), Inches(0.42)
        )
        title_frame = title_box.text_frame
        title_frame.text = titulo
        p = title_frame.paragraphs[0]
        p.font.size = Pt(27)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.LEFT

        return slide

    # Función para agregar logo UPC (copiar de la slide 1)
    def agregar_logo(slide):
        # Buscar el logo en la portada
        for shape in slide1.shapes:
            if shape.shape_type == 13:  # PICTURE
                if shape.left.inches > 8:  # Logo en esquina superior derecha
                    # Copiar el logo
                    el = shape.element
                    newel = copy.deepcopy(el)
                    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                    break

    # Función para agregar imagen de fondo decorativa
    def agregar_imagen_fondo(slide, left, top, width, height):
        # Crear un rectángulo como placeholder para la imagen
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = COLOR_GRIS_CLARO
        img_placeholder.line.fill.background()
        return img_placeholder

    # DIAPOSITIVA 2: Objetivo y Agenda
    slide2 = crear_slide_con_header("Objetivo y Agenda")
    agregar_logo(slide2)

    # Caja de objetivo
    obj_box = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.50), Inches(1.65),
        Inches(9.00), Inches(1.23)
    )
    obj_box.fill.solid()
    obj_box.fill.fore_color.rgb = COLOR_GRIS_CLARO
    obj_box.line.fill.background()

    # Línea vertical en la caja (casi invisible pero debe estar)
    line_vert = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.53), Inches(1.65),
        Inches(0.00), Inches(1.23)
    )
    line_vert.line.fill.background()

    # Título del objetivo
    obj_title = slide2.shapes.add_textbox(
        Inches(0.81), Inches(1.90),
        Inches(4.31), Inches(0.33)
    )
    tf = obj_title.text_frame
    tf.text = "🎯 Objetivo"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Texto del objetivo
    obj_text = slide2.shapes.add_textbox(
        Inches(0.81), Inches(2.30),
        Inches(8.61), Inches(0.29)
    )
    tf = obj_text.text_frame
    tf.word_wrap = True
    tf.text = "Enseñar la importancia de los datos para entrenar modelos de Inteligencia Artificial y cómo recopilar datos de calidad."
    p = tf.paragraphs[0]
    p.font.size = Pt(13.5)
    p.alignment = PP_ALIGN.LEFT

    # Título de agenda
    agenda_title = slide2.shapes.add_textbox(
        Inches(0.50), Inches(3.13),
        Inches(9.18), Inches(0.29)
    )
    tf = agenda_title.text_frame
    tf.text = "📋 Agenda"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Contenido de agenda
    agenda_text = slide2.shapes.add_textbox(
        Inches(0.50), Inches(3.56),
        Inches(9.00), Inches(1.29)
    )
    tf = agenda_text.text_frame
    tf.word_wrap = True

    items = [
        "Importancia de los datos en IA",
        "Tipos de datos para entrenar modelos",
        "Cómo recopilar datos de calidad",
        "Actividad práctica: recopilación de datos"
    ]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13.5)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 3: ¿Por qué son importantes los datos?
    slide3 = crear_slide_con_header("¿Por qué son importantes los datos?")
    agregar_logo(slide3)

    # Imagen de fondo decorativa
    agregar_imagen_fondo(slide3, 0.33, 1.64, 9.33, 1.55)

    # Texto introductorio
    intro_box = slide3.shapes.add_textbox(
        Inches(0.58), Inches(2.02),
        Inches(8.84), Inches(0.80)
    )
    tf = intro_box.text_frame
    tf.word_wrap = True
    tf.text = "Los datos son el combustible de la Inteligencia Artificial. Sin datos de calidad, los modelos de IA no pueden aprender correctamente."
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Tres cajas con conceptos clave
    conceptos = [
        ("💡", "Datos = Conocimiento", 0.33),
        ("🎯", "Calidad > Cantidad", 3.61),
        ("📊", "Variedad es clave", 6.89)
    ]

    for emoji, texto, left_pos in conceptos:
        # Caja de fondo
        caja = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(3.69),
            Inches(2.78), Inches(1.50)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_GRIS_CLARO
        caja.line.fill.background()

        # Emoji
        emoji_box = slide3.shapes.add_textbox(
            Inches(left_pos + 1.07), Inches(3.86),
            Inches(0.64), Inches(0.50)
        )
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

        # Texto - ajustar posición para centrarlo
        text_width = 1.78
        text_left = left_pos + (2.78 - text_width) / 2
        text_box = slide3.shapes.add_textbox(
            Inches(text_left), Inches(4.69),
            Inches(text_width), Inches(0.29)
        )
        tf = text_box.text_frame
        tf.text = texto
        p = tf.paragraphs[0]
        p.font.size = Pt(13.5)
        p.alignment = PP_ALIGN.CENTER

    # DIAPOSITIVA 4: Tipos de Datos para Entrenar IA
    slide4 = crear_slide_con_header("Tipos de Datos para Entrenar IA")
    agregar_logo(slide4)

    # Cuatro cajas con tipos de datos (usar BLANCO como en Sesión 5)
    tipos = [
        ("📸", "Imágenes", "Fotos de objetos, rostros,\nanimales, plantas", 0.25, 1.29),
        ("🔊", "Sonidos", "Voces, música,\nsonidos ambientales", 5.25, 1.29),
        ("🤸", "Poses", "Movimientos corporales,\nposturas, gestos", 0.25, 3.54),
        ("📝", "Texto", "Palabras, frases,\netiquetas, descripciones", 5.25, 3.54)
    ]

    for emoji, titulo, descripcion, left_pos, top_pos in tipos:
        # Caja de fondo BLANCA
        caja = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(2.00)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_BLANCO
        caja.line.fill.background()

        # Emoji y título en la misma línea
        header_box = slide4.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.27),
            Inches(4.04), Inches(0.29)
        )
        tf = header_box.text_frame
        tf.text = f"{emoji}  {titulo}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide4.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.73),
            Inches(4.04), Inches(0.50)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 5: ¿Cómo Recopilar Datos de Calidad?
    slide5 = crear_slide_con_header("¿Cómo Recopilar Datos de Calidad?")
    agregar_logo(slide5)

    # Imagen de fondo decorativa
    agregar_imagen_fondo(slide5, 0.60, 1.31, 8.81, 1.83)

    # Título grande centrado sobre la imagen
    title_grande = slide5.shapes.add_textbox(
        Inches(0.93), Inches(1.73),
        Inches(8.13), Inches(0.50)
    )
    tf = title_grande.text_frame
    tf.text = "5 Reglas de Oro"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Caja gris con las reglas
    reglas_box = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.40), Inches(3.64),
        Inches(7.20), Inches(1.55)
    )
    reglas_box.fill.solid()
    reglas_box.fill.fore_color.rgb = COLOR_GRIS_CLARO
    reglas_box.line.fill.background()

    # Línea vertical casi invisible
    line_vert2 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.43), Inches(3.64),
        Inches(0.00), Inches(1.55)
    )
    line_vert2.line.fill.background()

    # Título de la caja
    reglas_title = slide5.shapes.add_textbox(
        Inches(1.65), Inches(3.89),
        Inches(6.77), Inches(0.29)
    )
    tf = reglas_title.text_frame
    tf.text = "💭 Reglas clave"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Contenido de las reglas
    reglas_text = slide5.shapes.add_textbox(
        Inches(1.71), Inches(4.32),
        Inches(6.64), Inches(0.62)
    )
    tf = reglas_text.text_frame
    tf.word_wrap = True

    reglas = [
        "Cantidad suficiente • Variedad • Claridad • Balance • Relevancia"
    ]

    p = tf.paragraphs[0]
    p.text = reglas[0]
    p.font.size = Pt(13.5)
    p.alignment = PP_ALIGN.CENTER

    # DIAPOSITIVA 6: La IA ya está transformando el mundo
    slide6 = crear_slide_con_header("Ejemplos de Datos en Acción")
    agregar_logo(slide6)

    # Cuatro cajas BLANCAS con ejemplos
    ejemplos = [
        ("🐕", "Clasificador de Animales", "Reconoce perros, gatos y pájaros mediante fotos", 0.25, 1.71),
        ("🗣️", "Reconocedor de Voz", "Identifica diferentes voces y comandos hablados", 5.25, 1.71),
        ("👋", "Detector de Gestos", "Reconoce movimientos de manos para controlar dispositivos", 0.25, 3.54),
        ("📚", "Clasificador de Texto", "Categoriza comentarios como positivos o negativos", 5.25, 3.54)
    ]

    for emoji, titulo, descripcion, left_pos, top_pos in ejemplos:
        # Caja de fondo BLANCA
        caja = slide6.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(1.58)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_BLANCO
        caja.line.fill.background()

        # Emoji y título
        header_box = slide6.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.27),
            Inches(4.04), Inches(0.33)
        )
        tf = header_box.text_frame
        tf.text = f"{emoji}  {titulo}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide6.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.77),
            Inches(4.04), Inches(0.50)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 7: Pasos para Recopilar Datos
    slide7 = crear_slide_con_header("Pasos para Recopilar Datos")
    agregar_logo(slide7)

    # Cinco cajas BLANCAS con círculos naranjas numerados
    pasos = [
        ("1", "Definir categorías", "Decide qué clases tendrá tu modelo", 0.96),
        ("2", "Recopilar ejemplos", "Captura al menos 25-30 por categoría", 1.94),
        ("3", "Asegurar variedad", "Diferentes ángulos, iluminación, fondos", 2.92),
        ("4", "Verificar calidad", "Elimina datos borrosos o incorrectos", 3.91),
        ("5", "Organizar y etiquetar", "Clasifica los datos en grupos claros", 4.89)
    ]

    for numero, titulo, descripcion, top_pos in pasos:
        # Caja de fondo BLANCA
        caja = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.42), Inches(top_pos),
            Inches(8.98), Inches(0.63)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_BLANCO
        caja.line.fill.background()

        # Círculo naranja
        circulo = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.59), Inches(top_pos - 0.28),
            Inches(0.31), Inches(0.23)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = COLOR_NARANJA
        circulo.line.fill.background()

        # Número dentro del círculo
        num_box = slide7.shapes.add_textbox(
            Inches(0.71), Inches(top_pos - 0.26),
            Inches(0.08), Inches(0.15)
        )
        tf = num_box.text_frame
        tf.text = numero
        p = tf.paragraphs[0]
        p.font.size = Pt(10.8)
        p.alignment = PP_ALIGN.LEFT

        # Título del paso
        title_box = slide7.shapes.add_textbox(
            Inches(0.60), Inches(top_pos + 0.02),
            Inches(8.78), Inches(0.18)
        )
        tf = title_box.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide7.shapes.add_textbox(
            Inches(0.60), Inches(top_pos + 0.29),
            Inches(8.78), Inches(0.15)
        )
        tf = desc_box.text_frame
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(10.5)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 8: Herramientas para Recopilar Datos
    slide8 = crear_slide_con_header("🛠️ Herramientas y Recursos")
    agregar_logo(slide8)

    # Cuatro cajas BLANCAS con herramientas
    herramientas = [
        ("Teachable Machine", "Entrena modelos de IA sin código. Reconoce imágenes, sonidos y poses.", "teachablemachine.withgoogle.com", 0.25, 1.29),
        ("Scratch + ML", "Programa con bloques y añade extensiones de machine learning.", "scratch.mit.edu", 5.25, 1.29),
        ("Cámara y Micrófono", "Usa tu celular o computadora para capturar imágenes, videos y audio.", "Dispositivos personales", 0.25, 3.54),
        ("Colaboración", "Trabaja con compañeros para obtener más ejemplos y variedad.", "Trabajo en equipo", 5.25, 3.54)
    ]

    for titulo, descripcion, url, left_pos, top_pos in herramientas:
        # Caja de fondo BLANCA
        caja = slide8.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(2.00)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_BLANCO
        caja.line.fill.background()

        # Título
        title_box = slide8.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.27),
            Inches(4.04), Inches(0.29)
        )
        tf = title_box.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.size = Pt(15)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide8.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.73),
            Inches(4.04), Inches(0.50)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

        # URL
        url_box = slide8.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 1.48),
            Inches(4.04), Inches(0.21)
        )
        tf = url_box.text_frame
        tf.text = url
        p = tf.paragraphs[0]
        p.font.size = Pt(10.5)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 9: Actividad Práctica
    slide9 = crear_slide_con_header("👥 Actividad Práctica")
    agregar_logo(slide9)

    # Caja gris con tiempo
    tiempo_box = slide9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.33), Inches(1.43),
        Inches(4.42), Inches(0.83)
    )
    tiempo_box.fill.solid()
    tiempo_box.fill.fore_color.rgb = COLOR_GRIS_CLARO
    tiempo_box.line.fill.background()

    # Línea vertical
    line_tiempo = slide9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.36), Inches(1.43),
        Inches(0.00), Inches(0.83)
    )
    line_tiempo.line.fill.background()

    # Texto del tiempo
    tiempo_text = slide9.shapes.add_textbox(
        Inches(0.65), Inches(1.68),
        Inches(2.79), Inches(0.33)
    )
    tf = tiempo_text.text_frame
    tf.text = "⏱️ Tiempo: 20 minutos"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Título instrucciones
    instr_title = slide9.shapes.add_textbox(
        Inches(0.33), Inches(2.40),
        Inches(4.50), Inches(0.29)
    )
    tf = instr_title.text_frame
    tf.text = "📝 Instrucciones"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Lista de instrucciones
    instr_text = slide9.shapes.add_textbox(
        Inches(0.33), Inches(2.81),
        Inches(4.42), Inches(1.29)
    )
    tf = instr_text.text_frame
    tf.word_wrap = True

    instrucciones = [
        "Define tu proyecto de IA",
        "Identifica qué datos necesitas",
        "Determina cuántas categorías tendrás",
        "Comienza a recopilar ejemplos"
    ]

    for i, instruccion in enumerate(instrucciones):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = instruccion
        p.font.size = Pt(13.5)
        p.alignment = PP_ALIGN.LEFT

    # Título presentación
    pres_title = slide9.shapes.add_textbox(
        Inches(0.33), Inches(4.24),
        Inches(4.50), Inches(0.29)
    )
    tf = pres_title.text_frame
    tf.text = "💬 Presentación"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Texto presentación
    pres_text = slide9.shapes.add_textbox(
        Inches(0.33), Inches(4.70),
        Inches(4.50), Inches(0.25)
    )
    tf = pres_text.text_frame
    tf.text = "Comparte qué datos vas a recopilar y por qué."
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

    # Ficha de Proyecto (lado derecho)
    ficha_title = slide9.shapes.add_textbox(
        Inches(5.25), Inches(1.84),
        Inches(4.50), Inches(0.29)
    )
    tf = ficha_title.text_frame
    tf.text = "📋 Ficha de Recopilación"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Caja blanca con la ficha
    ficha_box = slide9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.25), Inches(2.26),
        Inches(4.42), Inches(2.32)
    )
    ficha_box.fill.solid()
    ficha_box.fill.fore_color.rgb = COLOR_BLANCO
    ficha_box.line.fill.background()

    # Campos de la ficha
    campos = [
        "Proyecto: _______________",
        "Tipo de datos: _______________",
        "Categorías: _______________",
        "Cantidad por categoría: _______________",
        "Herramienta: _______________",
        "Fecha límite: _______________"
    ]

    for i, campo in enumerate(campos):
        campo_box = slide9.shapes.add_textbox(
            Inches(5.47), Inches(2.48 + i * 0.33),
            Inches(4.05), Inches(0.21)
        )
        tf = campo_box.text_frame
        tf.text = campo
        p = tf.paragraphs[0]
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 10: Consejos y Errores Comunes
    slide10 = crear_slide_con_header("⚠️ Consejos y Errores Comunes")
    agregar_logo(slide10)

    # Errores comunes (cajas rojas) y buenas prácticas (cajas verdes)
    COLOR_ROJO_CLARO = RGBColor(254, 226, 226)
    COLOR_VERDE_CLARO = RGBColor(209, 250, 229)

    comparaciones = [
        ("❌ Pocos datos", "Sin suficientes ejemplos, la IA no aprende bien", COLOR_ROJO_CLARO, 0.25, 1.57),
        ("✅ Datos abundantes", "Al menos 25-30 ejemplos por categoría", COLOR_VERDE_CLARO, 5.25, 1.57),
        ("❌ Datos similares", "Sin variedad, el modelo no generaliza", COLOR_ROJO_CLARO, 0.25, 2.88),
        ("✅ Datos variados", "Diferentes ángulos, fondos y condiciones", COLOR_VERDE_CLARO, 5.25, 2.88),
        ("❌ Mala calidad", "Imágenes borrosas o audio con ruido", COLOR_ROJO_CLARO, 0.25, 4.20),
        ("✅ Alta calidad", "Datos claros, nítidos y representativos", COLOR_VERDE_CLARO, 5.25, 4.20)
    ]

    for titulo, descripcion, color, left_pos, top_pos in comparaciones:
        # Caja de fondo
        caja = slide10.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(1.07)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = color
        caja.line.fill.background()

        # Línea vertical
        line_caja = slide10.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos + 0.03), Inches(top_pos),
            Inches(0.00), Inches(1.07)
        )
        line_caja.line.fill.background()

        # Título
        title_box = slide10.shapes.add_textbox(
            Inches(left_pos + 0.26), Inches(top_pos + 0.20),
            Inches(4.12), Inches(0.29)
        )
        tf = title_box.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide10.shapes.add_textbox(
            Inches(left_pos + 0.26), Inches(top_pos + 0.61),
            Inches(4.12), Inches(0.25)
        )
        tf = desc_box.text_frame
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 11: Conclusión y Próximos Pasos
    slide11 = crear_slide_con_header("🎯 Conclusión y Próximos Pasos")
    agregar_logo(slide11)

    # Caja gris con puntos clave
    puntos_box = slide11.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.42), Inches(0.96),
        Inches(9.17), Inches(2.00)
    )
    puntos_box.fill.solid()
    puntos_box.fill.fore_color.rgb = COLOR_GRIS_CLARO
    puntos_box.line.fill.background()

    # Línea vertical
    line_puntos = slide11.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45), Inches(0.96),
        Inches(0.00), Inches(2.00)
    )
    line_puntos.line.fill.background()

    # Título puntos clave
    puntos_title = slide11.shapes.add_textbox(
        Inches(0.73), Inches(1.21),
        Inches(4.40), Inches(0.33)
    )
    tf = puntos_title.text_frame
    tf.text = "📌 Puntos Clave"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Lista de puntos clave
    puntos_text = slide11.shapes.add_textbox(
        Inches(0.73), Inches(1.75),
        Inches(8.60), Inches(0.96)
    )
    tf = puntos_text.text_frame
    tf.word_wrap = True

    puntos = [
        "Los datos son esenciales para entrenar modelos de IA",
        "La calidad y variedad son más importantes que la cantidad",
        "Organizar y etiquetar correctamente los datos es fundamental"
    ]

    for i, punto in enumerate(puntos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = punto
        p.font.size = Pt(13.5)
        p.alignment = PP_ALIGN.LEFT

    # Imagen de fondo para próxima sesión
    agregar_imagen_fondo(slide11, 0.42, 3.46, 9.17, 2.00)

    # Título próxima sesión
    prox_title = slide11.shapes.add_textbox(
        Inches(3.56), Inches(3.79),
        Inches(2.88), Inches(0.38)
    )
    tf = prox_title.text_frame
    tf.text = "🚀 Próxima Sesión"
    p = tf.paragraphs[0]
    p.font.size = Pt(22.5)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Texto próxima sesión
    prox_text = slide11.shapes.add_textbox(
        Inches(3.48), Inches(4.33),
        Inches(3.04), Inches(0.29)
    )
    tf = prox_text.text_frame
    tf.text = "Sesión 7: Entrenando el Modelo de IA"
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.alignment = PP_ALIGN.CENTER

    # Descripción
    prox_desc = slide11.shapes.add_textbox(
        Inches(2.25), Inches(4.83),
        Inches(5.51), Inches(0.25)
    )
    tf = prox_desc.text_frame
    tf.text = "Usaremos los datos recopilados para entrenar tu primer modelo de IA"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    # Guardar presentación
    prs.save('/home/user/MentesDigitales/Sesión 6.pptx')
    print("✅ Presentación 'Sesión 6.pptx' creada exitosamente siguiendo la estructura de Sesión 5!")

if __name__ == "__main__":
    crear_sesion_6()
