#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def crear_sesion_6():
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Definir colores del tema (basados en el PDF)
    COLOR_PRINCIPAL = RGBColor(41, 128, 185)  # Azul
    COLOR_ACENTO = RGBColor(52, 152, 219)  # Azul claro
    COLOR_TEXTO = RGBColor(44, 62, 80)  # Gris oscuro
    COLOR_ROJO = RGBColor(231, 76, 60)  # Rojo UPC
    COLOR_FONDO_CLARO = RGBColor(236, 240, 241)

    # Slide 1: Portada
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco

    # Fondo azul degradado superior
    shapes = slide1.shapes

    # Título principal
    title_box = shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Sesión 6:\nRecopilación y Preparación de Datos"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL
    title_para.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Mentes Digitales: Introducción a la IA con Google AI"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = COLOR_TEXTO
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Resumen de la sesión anterior
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Título
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "¿Qué hicimos en la Sesión 5?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    # Contenido
    content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Punto 1
    p1 = content_frame.paragraphs[0]
    p1.text = "✓ Identificamos problemas que la IA puede resolver"
    p1.font.size = Pt(24)
    p1.font.color.rgb = COLOR_TEXTO
    p1.space_after = Pt(20)

    # Punto 2
    p2 = content_frame.add_paragraph()
    p2.text = "✓ Lluvia de ideas para nuestros proyectos de IA"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_TEXTO
    p2.space_after = Pt(20)

    # Punto 3
    p3 = content_frame.add_paragraph()
    p3.text = "✓ Definimos el tipo de proyecto que queremos crear"
    p3.font.size = Pt(24)
    p3.font.color.rgb = COLOR_TEXTO
    p3.space_after = Pt(20)

    # Punto 4
    p4 = content_frame.add_paragraph()
    p4.text = "Ahora... ¿qué sigue?"
    p4.font.size = Pt(28)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_ROJO
    p4.alignment = PP_ALIGN.CENTER

    # Slide 3: Objetivo, Contenido y Actividad
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Título
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Sesión 6: Recopilación y Preparación de Datos"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL
    title_para.alignment = PP_ALIGN.CENTER

    # Columna 1: Objetivo
    obj_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.8), Inches(5))
    obj_frame = obj_box.text_frame
    obj_frame.word_wrap = True

    obj_title = obj_frame.paragraphs[0]
    obj_title.text = "Objetivo"
    obj_title.font.size = Pt(24)
    obj_title.font.bold = True
    obj_title.font.color.rgb = COLOR_PRINCIPAL
    obj_title.space_after = Pt(12)

    obj_content = obj_frame.add_paragraph()
    obj_content.text = "Enseñar la importancia de los datos para entrenar modelos de IA."
    obj_content.font.size = Pt(18)
    obj_content.font.color.rgb = COLOR_TEXTO
    obj_content.line_spacing = 1.3

    # Columna 2: Contenido
    cont_box = slide3.shapes.add_textbox(Inches(3.6), Inches(1.5), Inches(2.8), Inches(5))
    cont_frame = cont_box.text_frame
    cont_frame.word_wrap = True

    cont_title = cont_frame.paragraphs[0]
    cont_title.text = "Contenido"
    cont_title.font.size = Pt(24)
    cont_title.font.bold = True
    cont_title.font.color.rgb = COLOR_PRINCIPAL
    cont_title.space_after = Pt(12)

    cont_content = cont_frame.add_paragraph()
    cont_content.text = "Explicación de cómo recopilar y preparar datos para entrenar un modelo."
    cont_content.font.size = Pt(18)
    cont_content.font.color.rgb = COLOR_TEXTO
    cont_content.line_spacing = 1.3

    # Columna 3: Actividad
    act_box = slide3.shapes.add_textbox(Inches(6.7), Inches(1.5), Inches(2.8), Inches(5))
    act_frame = act_box.text_frame
    act_frame.word_wrap = True

    act_title = act_frame.paragraphs[0]
    act_title.text = "Actividad"
    act_title.font.size = Pt(24)
    act_title.font.bold = True
    act_title.font.color.rgb = COLOR_PRINCIPAL
    act_title.space_after = Pt(12)

    act_content = act_frame.add_paragraph()
    act_content.text = "Los estudiantes aprenderán a recopilar datos para su propio proyecto de IA (por ejemplo, imágenes de diferentes tipos de animales o sonidos)."
    act_content.font.size = Pt(18)
    act_content.font.color.rgb = COLOR_TEXTO
    act_content.line_spacing = 1.3

    # Slide 4: ¿Por qué son importantes los datos?
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "¿Por qué son importantes los datos?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    content_box = slide4.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p1 = content_frame.paragraphs[0]
    p1.text = "🧠 Los datos son el 'alimento' de la inteligencia artificial"
    p1.font.size = Pt(24)
    p1.font.color.rgb = COLOR_TEXTO
    p1.space_after = Pt(20)

    p2 = content_frame.add_paragraph()
    p2.text = "📚 Sin datos de calidad, la IA no puede aprender correctamente"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_TEXTO
    p2.space_after = Pt(20)

    p3 = content_frame.add_paragraph()
    p3.text = "🎯 Más datos = mejor precisión del modelo"
    p3.font.size = Pt(24)
    p3.font.color.rgb = COLOR_TEXTO
    p3.space_after = Pt(20)

    p4 = content_frame.add_paragraph()
    p4.text = "⚖️ La variedad de datos ayuda a la IA a generalizar mejor"
    p4.font.size = Pt(24)
    p4.font.color.rgb = COLOR_TEXTO

    # Slide 5: Tipos de datos para IA
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Tipos de Datos para Entrenar IA"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL
    title_para.alignment = PP_ALIGN.CENTER

    # Cuadro 1: Imágenes
    img_box = slide5.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(4), Inches(2.2))
    img_frame = img_box.text_frame
    img_frame.word_wrap = True

    img_title = img_frame.paragraphs[0]
    img_title.text = "📸 Imágenes"
    img_title.font.size = Pt(28)
    img_title.font.bold = True
    img_title.font.color.rgb = COLOR_PRINCIPAL
    img_title.space_after = Pt(12)

    img_content = img_frame.add_paragraph()
    img_content.text = "• Fotos de objetos\n• Expresiones faciales\n• Animales, plantas\n• Gestos y poses"
    img_content.font.size = Pt(18)
    img_content.font.color.rgb = COLOR_TEXTO
    img_content.line_spacing = 1.4

    # Cuadro 2: Sonidos
    sound_box = slide5.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4), Inches(2.2))
    sound_frame = sound_box.text_frame
    sound_frame.word_wrap = True

    sound_title = sound_frame.paragraphs[0]
    sound_title.text = "🔊 Sonidos"
    sound_title.font.size = Pt(28)
    sound_title.font.bold = True
    sound_title.font.color.rgb = COLOR_PRINCIPAL
    sound_title.space_after = Pt(12)

    sound_content = sound_frame.add_paragraph()
    sound_content.text = "• Voces humanas\n• Sonidos de animales\n• Instrumentos musicales\n• Sonidos ambientales"
    sound_content.font.size = Pt(18)
    sound_content.font.color.rgb = COLOR_TEXTO
    sound_content.line_spacing = 1.4

    # Cuadro 3: Poses
    pose_box = slide5.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(4), Inches(2.2))
    pose_frame = pose_box.text_frame
    pose_frame.word_wrap = True

    pose_title = pose_frame.paragraphs[0]
    pose_title.text = "🤸 Poses"
    pose_title.font.size = Pt(28)
    pose_title.font.bold = True
    pose_title.font.color.rgb = COLOR_PRINCIPAL
    pose_title.space_after = Pt(12)

    pose_content = pose_frame.add_paragraph()
    pose_content.text = "• Movimientos corporales\n• Posturas\n• Ejercicios\n• Bailes o gestos"
    pose_content.font.size = Pt(18)
    pose_content.font.color.rgb = COLOR_TEXTO
    pose_content.line_spacing = 1.4

    # Cuadro 4: Texto
    text_box = slide5.shapes.add_textbox(Inches(5.3), Inches(4.3), Inches(4), Inches(2.2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    text_title = text_frame.paragraphs[0]
    text_title.text = "📝 Texto"
    text_title.font.size = Pt(28)
    text_title.font.bold = True
    text_title.font.color.rgb = COLOR_PRINCIPAL
    text_title.space_after = Pt(12)

    text_content = text_frame.add_paragraph()
    text_content.text = "• Mensajes\n• Descripciones\n• Comentarios\n• Etiquetas"
    text_content.font.size = Pt(18)
    text_content.font.color.rgb = COLOR_TEXTO
    text_content.line_spacing = 1.4

    # Slide 6: ¿Cómo recopilar datos de calidad?
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "¿Cómo Recopilar Datos de Calidad?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    content_box = slide6.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    rules = [
        ("1️⃣ Cantidad suficiente", "Necesitas muchos ejemplos de cada categoría (mínimo 20-30)"),
        ("2️⃣ Variedad", "Incluye diferentes ángulos, iluminación, fondos, etc."),
        ("3️⃣ Claridad", "Asegúrate de que los datos sean claros y representativos"),
        ("4️⃣ Balance", "Ten una cantidad similar de datos para cada categoría"),
        ("5️⃣ Relevancia", "Los datos deben relacionarse directamente con tu proyecto")
    ]

    for i, (title, desc) in enumerate(rules):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = f"{title}: {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(15)
        p.line_spacing = 1.2

        # Hacer el título en negrita
        run = p.runs[0]
        run.font.bold = True

    # Slide 7: Ejemplo práctico
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Ejemplo: Clasificador de Animales"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    content_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p1 = content_frame.paragraphs[0]
    p1.text = "Proyecto: Una IA que identifique perros, gatos y pájaros"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ROJO
    p1.space_after = Pt(20)

    p2 = content_frame.add_paragraph()
    p2.text = "Datos necesarios:"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXTO
    p2.space_after = Pt(12)

    examples = [
        "✓ 30+ fotos de perros (diferentes razas, tamaños, colores)",
        "✓ 30+ fotos de gatos (diferentes poses, lugares)",
        "✓ 30+ fotos de pájaros (en vuelo, posados, diferentes especies)",
        "✓ Buena iluminación en todas las fotos",
        "✓ Fondos variados para cada categoría",
        "✓ Diferentes distancias: cerca y lejos"
    ]

    for example in examples:
        p = content_frame.add_paragraph()
        p.text = example
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(10)
        p.level = 1

    # Slide 8: Herramientas para recopilar datos
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Herramientas para Recopilar Datos"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    content_box = slide8.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    tools = [
        ("📱 Cámara de tu celular o tablet", "Para capturar imágenes y videos"),
        ("🎤 Micrófono", "Para grabar sonidos"),
        ("💻 Teachable Machine", "Permite capturar datos directamente desde el navegador"),
        ("🔍 Búsqueda de imágenes", "Con permiso, usar imágenes de internet (Creative Commons)"),
        ("👥 Colaboración", "Pedir ayuda a compañeros para tener más variedad")
    ]

    for i, (tool, desc) in enumerate(tools):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = f"{tool}\n   {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(18)
        p.line_spacing = 1.2

    # Slide 9: Preparación de datos
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Preparando tus Datos"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    content_box = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p1 = content_frame.paragraphs[0]
    p1.text = "Pasos para preparar tus datos:"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRINCIPAL
    p1.space_after = Pt(20)

    steps = [
        "1. Organiza tus datos por categorías (crea carpetas o clases)",
        "2. Elimina datos borrosos o de mala calidad",
        "3. Revisa que todas las categorías tengan cantidad similar",
        "4. Etiqueta correctamente cada grupo de datos",
        "5. Verifica que los datos representen bien lo que quieres enseñar"
    ]

    for step in steps:
        p = content_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(18)

    # Slide 10: Actividad práctica
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide10.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🎯 Actividad Práctica"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_ROJO
    title_para.alignment = PP_ALIGN.CENTER

    content_box = slide10.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p1 = content_frame.paragraphs[0]
    p1.text = "¡Es hora de recopilar datos para tu proyecto!"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRINCIPAL
    p1.space_after = Pt(20)
    p1.alignment = PP_ALIGN.CENTER

    p2 = content_frame.add_paragraph()
    p2.text = "Instrucciones:"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXTO
    p2.space_after = Pt(15)

    instructions = [
        "1. Revisa el proyecto que planeaste en la Sesión 5",
        "2. Define qué tipo de datos necesitas (imágenes, sonidos, poses)",
        "3. Determina cuántas categorías o clases tendrás",
        "4. Recopila al menos 25-30 ejemplos de cada categoría",
        "5. Organiza tus datos en carpetas o grupos",
        "6. Revisa la calidad y variedad de tus datos"
    ]

    for instruction in instructions:
        p = content_frame.add_paragraph()
        p.text = instruction
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(12)
        p.line_spacing = 1.2

    # Slide 11: Consejos finales
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide11.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "💡 Consejos Finales"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL

    content_box = slide11.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    tips = [
        "✨ Sé creativo: Usa diferentes ángulos, distancias y perspectivas",
        "🎨 Varía el contexto: Diferentes fondos y entornos",
        "☀️ Iluminación: Captura datos con buena luz",
        "🔄 Diversidad: Incluye diferentes variaciones de cada categoría",
        "🤝 Colabora: Trabaja con compañeros para obtener más datos",
        "⏰ Ten paciencia: La recopilación de datos lleva tiempo, ¡pero vale la pena!",
        "🎯 Calidad > Cantidad: Mejor pocos datos buenos que muchos malos"
    ]

    for i, tip in enumerate(tips):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = tip
        p.font.size = Pt(19)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(14)

    # Slide 12: Próxima sesión
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide12.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Próxima Sesión"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRINCIPAL
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide12.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    sub_p1 = subtitle_frame.paragraphs[0]
    sub_p1.text = "Sesión 7: Entrenando el Modelo de IA"
    sub_p1.font.size = Pt(32)
    sub_p1.font.color.rgb = COLOR_ACENTO
    sub_p1.alignment = PP_ALIGN.CENTER
    sub_p1.space_after = Pt(15)

    sub_p2 = subtitle_frame.add_paragraph()
    sub_p2.text = "¡Usaremos los datos que recopilaste hoy para entrenar tu modelo!"
    sub_p2.font.size = Pt(20)
    sub_p2.font.color.rgb = COLOR_TEXTO
    sub_p2.alignment = PP_ALIGN.CENTER

    # Guardar presentación
    prs.save('/home/user/MentesDigitales/Sesión 6.pptx')
    print("✅ Presentación 'Sesión 6.pptx' creada exitosamente!")

if __name__ == "__main__":
    crear_sesion_6()
