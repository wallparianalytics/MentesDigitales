#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Pt, Inches

def analizar_presentacion():
    prs = Presentation('/home/user/MentesDigitales/Sesión 5.pptx')

    print(f"Número de diapositivas: {len(prs.slides)}")
    print(f"Dimensiones: {prs.slide_width.inches} x {prs.slide_height.inches} pulgadas\n")

    for i, slide in enumerate(prs.slides):
        print(f"\n{'='*60}")
        print(f"DIAPOSITIVA {i+1}")
        print(f"{'='*60}")
        print(f"Layout: {slide.slide_layout.name}")

        # Analizar formas en la diapositiva
        for j, shape in enumerate(slide.shapes):
            print(f"\n  Forma {j+1}:")
            print(f"    Tipo: {shape.shape_type}")

            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(f"    Texto: {text[:100]}...")

            if hasattr(shape, "text_frame"):
                tf = shape.text_frame
                for k, paragraph in enumerate(tf.paragraphs):
                    if paragraph.text.strip():
                        print(f"      Párrafo {k+1}: {paragraph.text[:80]}")
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            if run.font.size:
                                print(f"        Tamaño fuente: {run.font.size.pt}")
                            if run.font.bold:
                                print(f"        Negrita: {run.font.bold}")
                            if run.font.color.type:
                                print(f"        Color tipo: {run.font.color.type}")
                            if paragraph.alignment:
                                print(f"        Alineación: {paragraph.alignment}")

            # Posición y tamaño
            if hasattr(shape, "left"):
                print(f"    Posición: ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
                print(f"    Tamaño: {shape.width.inches:.2f}\" x {shape.height.inches:.2f}\"")

            # Colores de relleno
            if hasattr(shape, "fill"):
                fill = shape.fill
                try:
                    if fill.type:
                        print(f"    Tipo de relleno: {fill.type}")
                        if fill.type == 1:  # SOLID
                            try:
                                rgb = fill.fore_color.rgb
                                print(f"    Color RGB: ({rgb[0]}, {rgb[1]}, {rgb[2]})")
                            except:
                                pass
                except:
                    pass

if __name__ == "__main__":
    analizar_presentacion()
