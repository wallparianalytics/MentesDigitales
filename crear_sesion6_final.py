#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import copy

def crear_sesion_6():
    # Cargar la presentación de la Sesión 5 como plantilla
    prs = Presentation('/home/user/MentesDigitales/Sesión 5.pptx')

    # Paleta de colores mejorada y variada
    COLOR_AZUL_HEADER = RGBColor(99, 102, 241)  # Azul índigo para headers
    COLOR_NARANJA = RGBColor(245, 158, 11)  # Naranja para acentos
    COLOR_MORADO = RGBColor(168, 85, 247)  # Morado para variedad
    COLOR_VERDE = RGBColor(34, 197, 94)  # Verde para positivo
    COLOR_ROSA = RGBColor(236, 72, 153)  # Rosa para destacar
    COLOR_CYAN = RGBColor(6, 182, 212)  # Cyan para contraste
    COLOR_GRIS_CLARO = RGBColor(241, 245, 249)  # Gris claro para cajas
    COLOR_BLANCO = RGBColor(255, 255, 255)  # Blanco
    COLOR_NEGRO = RGBColor(0, 0, 0)  # Negro para texto
    COLOR_TEXTO_OSCURO = RGBColor(30, 41, 59)  # Gris oscuro para texto principal
    COLOR_ROJO_CLARO = RGBColor(254, 226, 226)  # Rojo claro para errores
    COLOR_VERDE_CLARO = RGBColor(209, 250, 229)  # Verde claro para buenas prácticas

    # Mantener solo la primera diapositiva (portada) y eliminar el resto
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]

    # Modificar la portada (Diapositiva 1)
    slide1 = prs.slides[0]

    # Actualizar textos de la portada preservando el formato
    for shape in slide1.shapes:
        if hasattr(shape, "text_frame"):
            if "Sesión 5" in shape.text:
                # Preservar formato cambiando solo el texto del run
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "Sesión 5" in run.text:
                            run.text = "Sesión 6"
            elif "Prototipar tu Proyecto de IA" in shape.text:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = "Recopilación y Preparación de Datos"
            elif "Octubre 2024" in shape.text:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "Octubre 2024" in run.text:
                            run.text = "Noviembre 2024"

    # Función auxiliar para crear texto con formato
    def crear_texto_formateado(textbox, texto, fuente='Arial', tamano=12, negrita=False, color=None, alineacion=PP_ALIGN.LEFT):
        """Crea texto con formato aplicado correctamente a los runs"""
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = texto
        run.font.name = fuente
        run.font.size = Pt(tamano)
        if negrita:
            run.font.bold = True
        if color:
            run.font.color.rgb = color
        p.alignment = alineacion
        return p

    # Función auxiliar para crear una diapositiva con header azul
    def crear_slide_con_header(titulo):
        # Usar el primer layout disponible
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Eliminar todas las formas del layout
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)

        # Rectángulo azul superior
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.96)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_AZUL_HEADER
        header.line.fill.background()

        # Línea delgada debajo del header (casi invisible)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0.94),
            Inches(10), Inches(0.00)
        )
        line.line.fill.background()

        # Título en el header
        title_box = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.25),
            Inches(7), Inches(0.42)
        )
        crear_texto_formateado(title_box, titulo, 'Arial', 27, True, COLOR_BLANCO, PP_ALIGN.LEFT)

        return slide

    # Función para agregar logo UPC (copiar de la slide 1)
    def agregar_logo(slide):
        # Buscar el logo en la portada
        for shape in slide1.shapes:
            if shape.shape_type == 13:  # PICTURE
                if shape.left.inches > 8:  # Logo en esquina superior derecha
                    # Copiar el logo
                    el = shape.element
                    newel = copy.deepcopy(el)
                    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                    break

    # Función para crear iconos con formas geométricas
    def crear_icono_base_datos(slide, left, top, size, color):
        """Crea un icono de base de datos con cilindros"""
        # Cilindro superior
        cilindro1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top),
            Inches(size), Inches(size * 0.3)
        )
        cilindro1.fill.solid()
        cilindro1.fill.fore_color.rgb = color
        cilindro1.line.color.rgb = color
        cilindro1.line.width = Pt(2)

        # Cuerpo del cilindro
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top + size * 0.15),
            Inches(size), Inches(size * 0.5)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

        # Cilindro inferior
        cilindro2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top + size * 0.55),
            Inches(size), Inches(size * 0.3)
        )
        cilindro2.fill.solid()
        cilindro2.fill.fore_color.rgb = color
        cilindro2.line.color.rgb = color
        cilindro2.line.width = Pt(2)

    def crear_icono_cerebro(slide, left, top, size, color):
        """Crea un icono de cerebro estilizado"""
        # Círculo principal
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left), Inches(top),
            Inches(size), Inches(size)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = color
        circulo.line.color.rgb = COLOR_BLANCO
        circulo.line.width = Pt(3)

        # Líneas decorativas para simular neuronas
        for i in range(3):
            linea = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + size * 0.2 + i * 0.15), Inches(top + size * 0.3),
                Inches(0.02), Inches(size * 0.4)
            )
            linea.fill.solid()
            linea.fill.fore_color.rgb = COLOR_BLANCO
            linea.line.fill.background()

    def crear_icono_camara(slide, left, top, size, color):
        """Crea un icono de cámara"""
        # Cuerpo de la cámara
        cuerpo = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top + size * 0.2),
            Inches(size), Inches(size * 0.6)
        )
        cuerpo.fill.solid()
        cuerpo.fill.fore_color.rgb = color
        cuerpo.line.fill.background()

        # Lente (círculo blanco)
        lente = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + size * 0.3), Inches(top + size * 0.3),
            Inches(size * 0.4), Inches(size * 0.4)
        )
        lente.fill.solid()
        lente.fill.fore_color.rgb = COLOR_BLANCO
        lente.line.fill.background()

    def crear_icono_grafico(slide, left, top, size, color):
        """Crea un icono de gráfico de barras"""
        alturas = [0.3, 0.6, 0.4, 0.7]
        ancho_barra = size / 5

        for i, altura in enumerate(alturas):
            barra = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + i * ancho_barra + 0.05),
                Inches(top + size * (1 - altura)),
                Inches(ancho_barra - 0.05),
                Inches(size * altura)
            )
            barra.fill.solid()
            barra.fill.fore_color.rgb = color
            barra.line.fill.background()

    def crear_icono_checklist(slide, left, top, size, color):
        """Crea un icono de checklist"""
        # Fondo del documento
        doc = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(size), Inches(size * 1.2)
        )
        doc.fill.solid()
        doc.fill.fore_color.rgb = color
        doc.line.fill.background()

        # Marcas de verificación (checkmarks)
        for i in range(3):
            check = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + size * 0.2), Inches(top + 0.2 + i * 0.25),
                Inches(size * 0.6), Inches(0.04)
            )
            check.fill.solid()
            check.fill.fore_color.rgb = COLOR_BLANCO
            check.line.fill.background()

    # Función para agregar imagen de fondo decorativa con gradiente
    def agregar_imagen_fondo(slide, left, top, width, height, color):
        # Crear un rectángulo con color personalizado
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = color
        img_placeholder.line.fill.background()
        return img_placeholder

    # DIAPOSITIVA 2: Objetivo y Agenda
    slide2 = crear_slide_con_header("Objetivo y Agenda")
    agregar_logo(slide2)

    # Agregar icono de cerebro en la esquina
    crear_icono_cerebro(slide2, 8.8, 1.2, 0.6, COLOR_MORADO)

    # Caja de objetivo con gradiente morado
    obj_box = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.50), Inches(1.65),
        Inches(9.00), Inches(1.23)
    )
    obj_box.fill.solid()
    obj_box.fill.fore_color.rgb = RGBColor(237, 233, 254)  # Morado muy claro
    obj_box.line.fill.background()

    # Línea vertical en la caja (casi invisible pero debe estar)
    line_vert = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.53), Inches(1.65),
        Inches(0.00), Inches(1.23)
    )
    line_vert.line.fill.background()

    # Título del objetivo
    obj_title = slide2.shapes.add_textbox(
        Inches(0.81), Inches(1.90),
        Inches(4.31), Inches(0.33)
    )
    crear_texto_formateado(obj_title, "🎯 Objetivo", 'Arial', 18, True, COLOR_MORADO, PP_ALIGN.LEFT)

    # Texto del objetivo
    obj_text = slide2.shapes.add_textbox(
        Inches(0.81), Inches(2.30),
        Inches(8.61), Inches(0.29)
    )
    obj_text.text_frame.word_wrap = True
    crear_texto_formateado(obj_text, "Enseñar la importancia de los datos para entrenar modelos de Inteligencia Artificial y cómo recopilar datos de calidad.", 'Arial', 13.5, False, COLOR_TEXTO_OSCURO, PP_ALIGN.LEFT)

    # Título de agenda
    agenda_title = slide2.shapes.add_textbox(
        Inches(0.50), Inches(3.13),
        Inches(9.18), Inches(0.29)
    )
    crear_texto_formateado(agenda_title, "📋 Agenda", 'Arial', 15, True, COLOR_AZUL_HEADER, PP_ALIGN.LEFT)

    # Contenido de agenda
    agenda_text = slide2.shapes.add_textbox(
        Inches(0.50), Inches(3.56),
        Inches(9.00), Inches(1.29)
    )
    tf = agenda_text.text_frame
    tf.word_wrap = True

    items = [
        "Importancia de los datos en IA",
        "Tipos de datos para entrenar modelos",
        "Cómo recopilar datos de calidad",
        "Actividad práctica: recopilación de datos"
    ]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = 'Arial'
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 3: ¿Por qué son importantes los datos?
    slide3 = crear_slide_con_header("¿Por qué son importantes los datos?")
    agregar_logo(slide3)

    # Agregar icono de base de datos
    crear_icono_base_datos(slide3, 8.7, 1.1, 0.7, COLOR_AZUL_HEADER)

    # Imagen de fondo decorativa con azul más oscuro para contraste
    agregar_imagen_fondo(slide3, 0.33, 1.64, 9.33, 1.55, COLOR_AZUL_HEADER)

    # Texto introductorio en BLANCO sobre fondo azul oscuro
    intro_box = slide3.shapes.add_textbox(
        Inches(0.58), Inches(2.02),
        Inches(8.84), Inches(0.80)
    )
    tf = intro_box.text_frame
    tf.word_wrap = True
    tf.text = "Los datos son el combustible de la Inteligencia Artificial. Sin datos de calidad, los modelos de IA no pueden aprender correctamente."
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLANCO  # Ahora el blanco contrasta con el azul oscuro
    p.alignment = PP_ALIGN.CENTER

    # Tres cajas con conceptos clave en colores variados
    conceptos = [
        ("💡", "Datos = Conocimiento", 0.33, RGBColor(254, 243, 199)),  # Amarillo claro
        ("🎯", "Calidad > Cantidad", 3.61, RGBColor(219, 234, 254)),  # Azul claro
        ("📊", "Variedad es clave", 6.89, RGBColor(254, 226, 226))  # Rosa claro
    ]

    for emoji, texto, left_pos, color_fondo in conceptos:
        # Caja de fondo
        caja = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(3.69),
            Inches(2.78), Inches(1.50)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = color_fondo
        caja.line.fill.background()

        # Emoji
        emoji_box = slide3.shapes.add_textbox(
            Inches(left_pos + 1.07), Inches(3.86),
            Inches(0.64), Inches(0.50)
        )
        tf = emoji_box.text_frame
        tf.text = emoji
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.CENTER

        # Texto - ajustar posición para centrarlo
        text_width = 1.78
        text_left = left_pos + (2.78 - text_width) / 2
        text_box = slide3.shapes.add_textbox(
            Inches(text_left), Inches(4.69),
            Inches(text_width), Inches(0.29)
        )
        tf = text_box.text_frame
        tf.text = texto
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.CENTER

    # DIAPOSITIVA 4: Tipos de Datos para Entrenar IA
    slide4 = crear_slide_con_header("Tipos de Datos para Entrenar IA")
    agregar_logo(slide4)

    # Agregar icono de cámara
    crear_icono_camara(slide4, 8.8, 1.1, 0.7, COLOR_ROSA)

    # Cuatro cajas con tipos de datos con colores de borde
    tipos = [
        ("📸", "Imágenes", "Fotos de objetos, rostros,\nanimales, plantas", 0.25, 1.29, COLOR_ROSA),
        ("🔊", "Sonidos", "Voces, música,\nsonidos ambientales", 5.25, 1.29, COLOR_CYAN),
        ("🤸", "Poses", "Movimientos corporales,\nposturas, gestos", 0.25, 3.54, COLOR_NARANJA),
        ("📝", "Texto", "Palabras, frases,\netiquetas, descripciones", 5.25, 3.54, COLOR_VERDE)
    ]

    for emoji, titulo, descripcion, left_pos, top_pos, color_borde in tipos:
        # Caja de fondo BLANCA con borde de color
        caja = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(2.00)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_BLANCO
        # Agregar borde de color
        caja.line.color.rgb = color_borde
        caja.line.width = Pt(3)

        # Emoji y título en la misma línea
        header_box = slide4.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.27),
            Inches(4.04), Inches(0.29)
        )
        tf = header_box.text_frame
        tf.text = f"{emoji}  {titulo}"
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color_borde  # Color del borde también en el título
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide4.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.73),
            Inches(4.04), Inches(0.50)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 5: ¿Cómo Recopilar Datos de Calidad?
    slide5 = crear_slide_con_header("¿Cómo Recopilar Datos de Calidad?")
    agregar_logo(slide5)

    # Agregar icono de checklist
    crear_icono_checklist(slide5, 8.7, 1.0, 0.5, COLOR_VERDE)

    # Imagen de fondo decorativa con verde
    agregar_imagen_fondo(slide5, 0.60, 1.31, 8.81, 1.83, COLOR_VERDE)

    # Título grande centrado sobre la imagen en BLANCO
    title_grande = slide5.shapes.add_textbox(
        Inches(0.93), Inches(1.73),
        Inches(8.13), Inches(0.50)
    )
    tf = title_grande.text_frame
    tf.text = "5 Reglas de Oro"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER

    # Caja con gradiente amarillo claro para las reglas
    reglas_box = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.40), Inches(3.64),
        Inches(7.20), Inches(1.55)
    )
    reglas_box.fill.solid()
    reglas_box.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Amarillo claro
    reglas_box.line.fill.background()

    # Línea vertical casi invisible
    line_vert2 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.43), Inches(3.64),
        Inches(0.00), Inches(1.55)
    )
    line_vert2.line.fill.background()

    # Título de la caja
    reglas_title = slide5.shapes.add_textbox(
        Inches(1.65), Inches(3.89),
        Inches(6.77), Inches(0.29)
    )
    tf = reglas_title.text_frame
    tf.text = "💭 Reglas clave"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.CENTER

    # Contenido de las reglas
    reglas_text = slide5.shapes.add_textbox(
        Inches(1.71), Inches(4.32),
        Inches(6.64), Inches(0.62)
    )
    tf = reglas_text.text_frame
    tf.word_wrap = True

    reglas = [
        "Cantidad suficiente • Variedad • Claridad • Balance • Relevancia"
    ]

    p = tf.paragraphs[0]
    p.text = reglas[0]
    p.font.name = 'Arial'
    p.font.size = Pt(13.5)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.CENTER

    # DIAPOSITIVA 6: La IA ya está transformando el mundo
    slide6 = crear_slide_con_header("Ejemplos de Datos en Acción")
    agregar_logo(slide6)

    # Agregar icono de gráfico
    crear_icono_grafico(slide6, 8.7, 1.1, 0.7, COLOR_NARANJA)

    # Cuatro cajas con ejemplos - colores de fondo variados
    ejemplos = [
        ("🐕", "Clasificador de Animales", "Reconoce perros, gatos y pájaros mediante fotos", 0.25, 1.71, RGBColor(254, 243, 199)),
        ("🗣️", "Reconocedor de Voz", "Identifica diferentes voces y comandos hablados", 5.25, 1.71, RGBColor(219, 234, 254)),
        ("👋", "Detector de Gestos", "Reconoce movimientos de manos para controlar dispositivos", 0.25, 3.54, RGBColor(254, 226, 226)),
        ("📚", "Clasificador de Texto", "Categoriza comentarios como positivos o negativos", 5.25, 3.54, RGBColor(220, 252, 231))
    ]

    for emoji, titulo, descripcion, left_pos, top_pos, color_fondo in ejemplos:
        # Caja de fondo con color
        caja = slide6.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(1.58)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = color_fondo
        caja.line.fill.background()

        # Emoji y título
        header_box = slide6.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.27),
            Inches(4.04), Inches(0.33)
        )
        tf = header_box.text_frame
        tf.text = f"{emoji}  {titulo}"
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide6.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.77),
            Inches(4.04), Inches(0.50)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 7: Pasos para Recopilar Datos - REDISEÑADA CON TARJETAS
    slide7 = crear_slide_con_header("Pasos para Recopilar Datos")
    agregar_logo(slide7)

    # Rediseñar con tarjetas más grandes y mejor espaciadas
    pasos = [
        ("1", "Definir categorías", "Decide qué clases tendrá tu modelo", 0.33, 1.25, COLOR_ROSA),
        ("2", "Recopilar ejemplos", "Captura al menos 25-30 por categoría", 3.42, 1.25, COLOR_NARANJA),
        ("3", "Asegurar variedad", "Diferentes ángulos, iluminación, fondos", 6.50, 1.25, COLOR_CYAN),
        ("4", "Verificar calidad", "Elimina datos borrosos o incorrectos", 0.33, 3.40, COLOR_MORADO),
        ("5", "Organizar y etiquetar", "Clasifica los datos en grupos claros", 3.42, 3.40, COLOR_VERDE)
    ]

    for numero, titulo, descripcion, left_pos, top_pos, color_acento in pasos:
        # Caja de fondo BLANCA con espacio adecuado
        caja = slide7.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(2.92), Inches(1.90)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = COLOR_BLANCO
        caja.line.color.rgb = color_acento
        caja.line.width = Pt(2)

        # Círculo de color con número DENTRO de la caja
        circulo = slide7.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left_pos + 0.20), Inches(top_pos + 0.20),
            Inches(0.50), Inches(0.50)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = color_acento
        circulo.line.fill.background()

        # Número dentro del círculo - centrado
        num_box = slide7.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 0.27),
            Inches(0.50), Inches(0.35)
        )
        tf = num_box.text_frame
        tf.text = numero
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER

        # Título del paso
        title_box = slide7.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 0.85),
            Inches(2.52), Inches(0.40)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color_acento
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide7.shapes.add_textbox(
            Inches(left_pos + 0.20), Inches(top_pos + 1.30),
            Inches(2.52), Inches(0.45)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 8: Herramientas para Recopilar Datos
    slide8 = crear_slide_con_header("🛠️ Herramientas y Recursos")
    agregar_logo(slide8)

    # Cuatro cajas con herramientas - fondos de colores variados
    herramientas = [
        ("Teachable Machine", "Entrena modelos de IA sin código. Reconoce imágenes, sonidos y poses.", "teachablemachine.withgoogle.com", 0.25, 1.29, RGBColor(254, 243, 199)),
        ("Scratch + ML", "Programa con bloques y añade extensiones de machine learning.", "scratch.mit.edu", 5.25, 1.29, RGBColor(237, 233, 254)),
        ("Cámara y Micrófono", "Usa tu celular o computadora para capturar imágenes, videos y audio.", "Dispositivos personales", 0.25, 3.54, RGBColor(219, 234, 254)),
        ("Colaboración", "Trabaja con compañeros para obtener más ejemplos y variedad.", "Trabajo en equipo", 5.25, 3.54, RGBColor(220, 252, 231))
    ]

    for titulo, descripcion, url, left_pos, top_pos, color_fondo in herramientas:
        # Caja de fondo con color
        caja = slide8.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(2.00)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = color_fondo
        caja.line.fill.background()

        # Título
        title_box = slide8.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.27),
            Inches(4.04), Inches(0.29)
        )
        tf = title_box.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide8.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 0.73),
            Inches(4.04), Inches(0.50)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

        # URL
        url_box = slide8.shapes.add_textbox(
            Inches(left_pos + 0.27), Inches(top_pos + 1.48),
            Inches(4.04), Inches(0.21)
        )
        tf = url_box.text_frame
        tf.text = url
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 9: EJERCICIO PARA LOS ALUMNOS
    slide9 = crear_slide_con_header("📝 Ejercicio: Mi Primer Dataset")
    agregar_logo(slide9)

    # Caja principal con instrucciones del ejercicio
    ejercicio_box = slide9.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.33), Inches(1.20),
        Inches(9.33), Inches(4.30)
    )
    ejercicio_box.fill.solid()
    ejercicio_box.fill.fore_color.rgb = RGBColor(243, 244, 246)  # Gris muy claro
    ejercicio_box.line.color.rgb = COLOR_AZUL_HEADER
    ejercicio_box.line.width = Pt(3)

    # Título del ejercicio
    ejercicio_title = slide9.shapes.add_textbox(
        Inches(0.60), Inches(1.50),
        Inches(8.80), Inches(0.40)
    )
    tf = ejercicio_title.text_frame
    tf.text = "🎓 Proyecto: Construye un Clasificador de Objetos Escolares"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZUL_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Descripción del ejercicio
    desc_ejercicio = slide9.shapes.add_textbox(
        Inches(0.80), Inches(2.00),
        Inches(8.40), Inches(0.50)
    )
    tf = desc_ejercicio.text_frame
    tf.word_wrap = True
    tf.text = "Crea un modelo de IA que pueda reconocer 3 objetos escolares diferentes usando la cámara. Vas a recopilar tus propios datos y entrenar el modelo."
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.LEFT

    # Sección 1: Paso 1
    paso1_box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.80), Inches(2.65),
        Inches(4.00), Inches(1.25)
    )
    paso1_box.fill.solid()
    paso1_box.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Amarillo claro
    paso1_box.line.fill.background()

    paso1_title = slide9.shapes.add_textbox(
        Inches(1.00), Inches(2.80),
        Inches(3.60), Inches(0.30)
    )
    tf = paso1_title.text_frame
    tf.text = "1️⃣ Selecciona tus objetos"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_NARANJA
    p.alignment = PP_ALIGN.LEFT

    paso1_text = slide9.shapes.add_textbox(
        Inches(1.00), Inches(3.15),
        Inches(3.60), Inches(0.60)
    )
    tf = paso1_text.text_frame
    tf.word_wrap = True
    tf.text = "• Escoge 3 objetos: lápiz, libro, calculadora\n• Asegúrate de que sean fáciles de distinguir\n• Ten los objetos a la mano"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.LEFT

    # Sección 2: Paso 2
    paso2_box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.05), Inches(2.65),
        Inches(4.00), Inches(1.25)
    )
    paso2_box.fill.solid()
    paso2_box.fill.fore_color.rgb = RGBColor(219, 234, 254)  # Azul claro
    paso2_box.line.fill.background()

    paso2_title = slide9.shapes.add_textbox(
        Inches(5.25), Inches(2.80),
        Inches(3.60), Inches(0.30)
    )
    tf = paso2_title.text_frame
    tf.text = "2️⃣ Recopila las imágenes"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.LEFT

    paso2_text = slide9.shapes.add_textbox(
        Inches(5.25), Inches(3.15),
        Inches(3.60), Inches(0.60)
    )
    tf = paso2_text.text_frame
    tf.word_wrap = True
    tf.text = "• 30 fotos de cada objeto (90 en total)\n• Varía: ángulos, distancia, iluminación\n• Usa fondos diferentes"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.LEFT

    # Sección 3: Paso 3
    paso3_box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.80), Inches(4.05),
        Inches(4.00), Inches(1.25)
    )
    paso3_box.fill.solid()
    paso3_box.fill.fore_color.rgb = RGBColor(220, 252, 231)  # Verde claro
    paso3_box.line.fill.background()

    paso3_title = slide9.shapes.add_textbox(
        Inches(1.00), Inches(4.20),
        Inches(3.60), Inches(0.30)
    )
    tf = paso3_title.text_frame
    tf.text = "3️⃣ Entrena tu modelo"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.LEFT

    paso3_text = slide9.shapes.add_textbox(
        Inches(1.00), Inches(4.55),
        Inches(3.60), Inches(0.60)
    )
    tf = paso3_text.text_frame
    tf.word_wrap = True
    tf.text = "• Ve a teachablemachine.withgoogle.com\n• Sube tus fotos en 3 clases\n• Entrena el modelo y pruébalo"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.LEFT

    # Sección 4: Paso 4
    paso4_box = slide9.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.05), Inches(4.05),
        Inches(4.00), Inches(1.25)
    )
    paso4_box.fill.solid()
    paso4_box.fill.fore_color.rgb = RGBColor(254, 226, 226)  # Rosa claro
    paso4_box.line.fill.background()

    paso4_title = slide9.shapes.add_textbox(
        Inches(5.25), Inches(4.20),
        Inches(3.60), Inches(0.30)
    )
    tf = paso4_title.text_frame
    tf.text = "4️⃣ Presenta tu proyecto"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSA
    p.alignment = PP_ALIGN.LEFT

    paso4_text = slide9.shapes.add_textbox(
        Inches(5.25), Inches(4.55),
        Inches(3.60), Inches(0.60)
    )
    tf = paso4_text.text_frame
    tf.word_wrap = True
    tf.text = "• Muestra tu modelo funcionando\n• Explica qué objetos reconoce\n• Comparte qué aprendiste"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 10: Actividad Práctica (la original simplificada)
    slide10 = crear_slide_con_header("👥 Actividad Práctica en Clase")
    agregar_logo(slide10)

    # Caja con gradiente cyan para tiempo
    tiempo_box = slide10.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.33), Inches(1.43),
        Inches(4.42), Inches(0.83)
    )
    tiempo_box.fill.solid()
    tiempo_box.fill.fore_color.rgb = RGBColor(207, 250, 254)  # Cyan muy claro
    tiempo_box.line.fill.background()

    # Línea vertical
    line_tiempo = slide10.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.36), Inches(1.43),
        Inches(0.00), Inches(0.83)
    )
    line_tiempo.line.fill.background()

    # Texto del tiempo
    tiempo_text = slide10.shapes.add_textbox(
        Inches(0.65), Inches(1.68),
        Inches(2.79), Inches(0.33)
    )
    tf = tiempo_text.text_frame
    tf.text = "⏱️ Tiempo: 20 minutos"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.LEFT

    # Título instrucciones
    instr_title = slide10.shapes.add_textbox(
        Inches(0.33), Inches(2.40),
        Inches(4.50), Inches(0.29)
    )
    tf = instr_title.text_frame
    tf.text = "📝 Instrucciones"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_MORADO
    p.alignment = PP_ALIGN.LEFT

    # Lista de instrucciones
    instr_text = slide10.shapes.add_textbox(
        Inches(0.33), Inches(2.81),
        Inches(4.42), Inches(1.29)
    )
    tf = instr_text.text_frame
    tf.word_wrap = True

    instrucciones = [
        "Define tu proyecto de IA",
        "Identifica qué datos necesitas",
        "Determina cuántas categorías tendrás",
        "Comienza a recopilar ejemplos"
    ]

    for i, instruccion in enumerate(instrucciones):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = instruccion
        p.font.name = 'Arial'
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # Título presentación
    pres_title = slide10.shapes.add_textbox(
        Inches(0.33), Inches(4.24),
        Inches(4.50), Inches(0.29)
    )
    tf = pres_title.text_frame
    tf.text = "💬 Presentación"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_NARANJA
    p.alignment = PP_ALIGN.LEFT

    # Texto presentación
    pres_text = slide10.shapes.add_textbox(
        Inches(0.33), Inches(4.70),
        Inches(4.50), Inches(0.25)
    )
    tf = pres_text.text_frame
    tf.text = "Comparte qué datos vas a recopilar y por qué."
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.LEFT

    # Ficha de Proyecto (lado derecho) con gradiente rosa
    ficha_title = slide10.shapes.add_textbox(
        Inches(5.25), Inches(1.84),
        Inches(4.50), Inches(0.29)
    )
    tf = ficha_title.text_frame
    tf.text = "📋 Ficha de Recopilación"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROSA
    p.alignment = PP_ALIGN.LEFT

    # Caja con fondo amarillo claro
    ficha_box = slide10.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.25), Inches(2.26),
        Inches(4.42), Inches(2.32)
    )
    ficha_box.fill.solid()
    ficha_box.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Amarillo claro
    ficha_box.line.fill.background()

    # Campos de la ficha
    campos = [
        "Proyecto: _______________",
        "Tipo de datos: _______________",
        "Categorías: _______________",
        "Cantidad por categoría: _______________",
        "Herramienta: _______________",
        "Fecha límite: _______________"
    ]

    for i, campo in enumerate(campos):
        campo_box = slide10.shapes.add_textbox(
            Inches(5.47), Inches(2.48 + i * 0.33),
            Inches(4.05), Inches(0.21)
        )
        tf = campo_box.text_frame
        tf.text = campo
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 11: Consejos y Errores Comunes
    slide11 = crear_slide_con_header("⚠️ Consejos y Errores Comunes")
    agregar_logo(slide11)

    # Errores comunes (cajas rojas) y buenas prácticas (cajas verdes)
    comparaciones = [
        ("❌ Pocos datos", "Sin suficientes ejemplos, la IA no aprende bien", COLOR_ROJO_CLARO, 0.25, 1.57),
        ("✅ Datos abundantes", "Al menos 25-30 ejemplos por categoría", COLOR_VERDE_CLARO, 5.25, 1.57),
        ("❌ Datos similares", "Sin variedad, el modelo no generaliza", COLOR_ROJO_CLARO, 0.25, 2.88),
        ("✅ Datos variados", "Diferentes ángulos, fondos y condiciones", COLOR_VERDE_CLARO, 5.25, 2.88),
        ("❌ Mala calidad", "Imágenes borrosas o audio con ruido", COLOR_ROJO_CLARO, 0.25, 4.20),
        ("✅ Alta calidad", "Datos claros, nítidos y representativos", COLOR_VERDE_CLARO, 5.25, 4.20)
    ]

    for titulo, descripcion, color, left_pos, top_pos in comparaciones:
        # Caja de fondo
        caja = slide11.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos), Inches(top_pos),
            Inches(4.50), Inches(1.07)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = color
        caja.line.fill.background()

        # Línea vertical
        line_caja = slide11.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_pos + 0.03), Inches(top_pos),
            Inches(0.00), Inches(1.07)
        )
        line_caja.line.fill.background()

        # Título
        title_box = slide11.shapes.add_textbox(
            Inches(left_pos + 0.26), Inches(top_pos + 0.20),
            Inches(4.12), Inches(0.29)
        )
        tf = title_box.text_frame
        tf.text = titulo
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

        # Descripción
        desc_box = slide11.shapes.add_textbox(
            Inches(left_pos + 0.26), Inches(top_pos + 0.61),
            Inches(4.12), Inches(0.25)
        )
        tf = desc_box.text_frame
        tf.text = descripcion
        p = tf.paragraphs[0]
        p.font.name = 'Arial'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # DIAPOSITIVA 12: Conclusión y Próximos Pasos
    slide12 = crear_slide_con_header("🎯 Conclusión y Próximos Pasos")
    agregar_logo(slide12)

    # Caja con gradiente azul claro para puntos clave
    puntos_box = slide12.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.42), Inches(0.96),
        Inches(9.17), Inches(2.00)
    )
    puntos_box.fill.solid()
    puntos_box.fill.fore_color.rgb = RGBColor(219, 234, 254)  # Azul claro
    puntos_box.line.fill.background()

    # Línea vertical
    line_puntos = slide12.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45), Inches(0.96),
        Inches(0.00), Inches(2.00)
    )
    line_puntos.line.fill.background()

    # Título puntos clave
    puntos_title = slide12.shapes.add_textbox(
        Inches(0.73), Inches(1.21),
        Inches(4.40), Inches(0.33)
    )
    tf = puntos_title.text_frame
    tf.text = "📌 Puntos Clave"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZUL_HEADER
    p.alignment = PP_ALIGN.LEFT

    # Lista de puntos clave
    puntos_text = slide12.shapes.add_textbox(
        Inches(0.73), Inches(1.75),
        Inches(8.60), Inches(0.96)
    )
    tf = puntos_text.text_frame
    tf.word_wrap = True

    puntos = [
        "Los datos son esenciales para entrenar modelos de IA",
        "La calidad y variedad son más importantes que la cantidad",
        "Organizar y etiquetar correctamente los datos es fundamental"
    ]

    for i, punto in enumerate(puntos):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = punto
        p.font.name = 'Arial'
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT

    # Imagen de fondo para próxima sesión con morado
    agregar_imagen_fondo(slide12, 0.42, 3.46, 9.17, 2.00, COLOR_MORADO)

    # Título próxima sesión en BLANCO
    prox_title = slide12.shapes.add_textbox(
        Inches(3.56), Inches(3.79),
        Inches(2.88), Inches(0.38)
    )
    tf = prox_title.text_frame
    tf.text = "🚀 Próxima Sesión"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(22.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER

    # Texto próxima sesión en BLANCO
    prox_text = slide12.shapes.add_textbox(
        Inches(3.48), Inches(4.33),
        Inches(3.04), Inches(0.29)
    )
    tf = prox_text.text_frame
    tf.text = "Sesión 7: Entrenando el Modelo de IA"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER

    # Descripción en BLANCO
    prox_desc = slide12.shapes.add_textbox(
        Inches(2.25), Inches(4.83),
        Inches(5.51), Inches(0.25)
    )
    tf = prox_desc.text_frame
    tf.text = "Usaremos los datos recopilados para entrenar tu primer modelo de IA"
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER

    # Guardar presentación
    prs.save('/home/user/MentesDigitales/Sesión 6.pptx')
    print("✅ Presentación 'Sesión 6.pptx' creada exitosamente con gráficos, iconos y ejercicio para alumnos!")

if __name__ == "__main__":
    crear_sesion_6()
